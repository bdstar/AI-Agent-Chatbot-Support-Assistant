import os
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_chroma import Chroma
from app.config import settings
from app.utils.helpers import get_file_extension, is_supported_document


def load_documents_from_directory(directory: str = None) -> list:
    """
    Load all supported documents from the Documents directory.
    Supports: PDF, TXT, PPTX, and images.
    Returns a list of LangChain Document objects.
    """
    if directory is None:
        directory = settings.DOCUMENTS_DIR

    if not os.path.exists(directory):
        os.makedirs(directory, exist_ok=True)
        print(f"[INFO] Created Documents directory: {directory}")
        return []

    all_documents = []
    files_processed = 0

    for root, dirs, files in os.walk(directory):
        for filename in files:
            if not is_supported_document(filename):
                continue

            filepath = os.path.join(root, filename)
            ext = get_file_extension(filename)

            try:
                if ext == ".pdf":
                    loader = PyPDFLoader(filepath)
                    docs = loader.load()
                elif ext in (".txt", ".md"):
                    loader = TextLoader(filepath, encoding="utf-8")
                    docs = loader.load()
                elif ext in (".pptx", ".ppt"):
                    # Use a simple text extraction for PPT files
                    docs = _load_pptx(filepath)
                elif ext in (".png", ".jpg", ".jpeg"):
                    # Skip images for now (would require OCR)
                    print(f"[SKIP] Skipping image file (OCR not configured): {filename}")
                    continue
                else:
                    continue

                # Add source metadata
                for doc in docs:
                    doc.metadata["source"] = filepath
                    doc.metadata["filename"] = filename

                all_documents.extend(docs)
                files_processed += 1
                print(f"[OK] Loaded: {filename} ({len(docs)} pages/sections)")

            except Exception as e:
                print(f"[ERROR] Error loading {filename}: {e}")

    print(f"\n[INFO] Total: {files_processed} files loaded, {len(all_documents)} document sections")
    return all_documents


def _load_pptx(filepath: str) -> list:
    """Load a PowerPoint file and extract text content."""
    from pptx import Presentation
    from langchain_core.documents import Document

    prs = Presentation(filepath)
    documents = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)

        if text_parts:
            content = "\n".join(text_parts)
            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": filepath,
                        "page": slide_num,
                        "type": "pptx",
                    },
                )
            )

    return documents


def split_documents(documents: list) -> list:
    """
    Split documents into smaller chunks for embedding.
    Uses RecursiveCharacterTextSplitter for intelligent chunking.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = text_splitter.split_documents(documents)
    print(f"[INFO] Split into {len(chunks)} chunks (chunk_size={settings.CHUNK_SIZE}, overlap={settings.CHUNK_OVERLAP})")
    return chunks


def create_vector_store(chunks: list) -> Chroma:
    """
    Create embeddings and store them in ChromaDB.
    Returns the Chroma vector store instance.
    """
    embeddings = OpenAIEmbeddings(
        model=settings.EMBEDDING_MODEL,
        openai_api_key=settings.OPENAI_API_KEY,
    )

    vectorstore = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=settings.CHROMA_DIR,
        collection_name="life_insurance_docs",
    )

    print(f"[OK] Vector store created at: {settings.CHROMA_DIR}")
    return vectorstore


def ingest_documents() -> dict:
    """
    Full ingestion pipeline:
    1. Load documents from directory
    2. Split into chunks
    3. Create embeddings and store in ChromaDB

    Returns a summary dict.
    """
    print("\n[START] Starting document ingestion pipeline...")
    print("=" * 50)

    # Step 1: Load documents
    documents = load_documents_from_directory()
    if not documents:
        return {
            "status": "warning",
            "documents_processed": 0,
            "chunks_created": 0,
            "message": "No documents found in the Documents directory. Please add PDF, TXT, or PPTX files.",
        }

    # Step 2: Split into chunks
    chunks = split_documents(documents)

    # Step 3: Create vector store
    create_vector_store(chunks)

    print("=" * 50)
    print("[OK] Ingestion complete!\n")

    return {
        "status": "success",
        "documents_processed": len(documents),
        "chunks_created": len(chunks),
        "message": f"Successfully processed {len(documents)} document sections into {len(chunks)} chunks.",
    }
